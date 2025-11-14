from flask import Flask, render_template, request, send_file, g
import os
import tempfile
from pdf2docx import Converter
from pptx import Presentation
from PIL import Image
import fitz  # PyMuPDF for PDF rendering

app = Flask(__name__)

# --- CONFIGURATION ---
UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 32 * 1024 * 1024


# ---------- Helper Functions ---------- #

def pdf_to_docx(pdf_path, output_path):
    """Convert PDF -> DOCX with layout retention."""
    try:
        cv = Converter(pdf_path)
        cv.convert(output_path, start=0, end=None)
        cv.close()
        return "Converted successfully"
    except Exception as e:
        return f"Conversion error: {str(e)}"


def pdf_to_ppt(pdf_path, output_pptx):
    """Convert PDF -> PPTX by adding each page as an image."""
    try:
        doc = fitz.open(pdf_path)
        prs = Presentation()
        blank_slide = prs.slide_layouts[6]

        with tempfile.TemporaryDirectory() as temp_dir:
            for page in doc:
                pix = page.get_pixmap()
                img_path = os.path.join(temp_dir, f"page_{page.number}.png")
                pix.save(img_path)
                
                slide = prs.slides.add_slide(blank_slide)
                slide.shapes.add_picture(img_path, 0, 0, prs.slide_width, prs.slide_height)
                
        prs.save(output_pptx)
        return "Converted successfully"
    except Exception as e:
        return f"Conversion error: {str(e)}"


def image_to_pdf(input_path, output_path):
    """Convert JPG, PNG, WEBP etc. -> PDF."""
    try:
        img = Image.open(input_path).convert('RGB')
        img.save(output_path)
        return "Converted successfully"
    except Exception as e:
        return f"Conversion error: {str(e)}"

# --- NEW FUNCTION ---
def image_to_image(input_path, output_path, target_format):
    """Convert between image types (JPG, PNG, WEBP)."""
    try:
        img = Image.open(input_path)
        # Handle formats that need RGB (like JPG)
        if target_format.lower() == 'jpg' or target_format.lower() == 'jpeg':
             img = img.convert('RGB')
        img.save(output_path)
        return "Converted successfully"
    except Exception as e:
        return f"Conversion error: {str(e)}"

# --- NEW FUNCTION ---
def pdf_to_txt(input_path, output_path):
    """Extract text from PDF -> TXT."""
    try:
        doc = fitz.open(input_path)
        with open(output_path, 'w', encoding='utf-8') as txt_file:
            for page in doc:
                txt_file.write(page.get_text())
                txt_file.write("\n")
        return "Converted successfully"
    except Exception as e:
        return f"Conversion error: {str(e)}"


# ---------- Flask Routes ---------- #

@app.route('/')
def index():
    """Serves the main index.html page."""
    return render_template('index.html')


@app.route('/convert', methods=['POST'])
def convert_file():
    """Handles the file upload and conversion logic."""
    
    file = request.files.get('file')
    target_format = request.form.get('target_format') # e.g., "docx"

    if not file or not target_format:
        return "<h3>Missing file or target format!</h3><a href='/'>Go Back</a>", 400

    original_name, original_ext = os.path.splitext(file.filename)
    source_ext = original_ext.lower().replace('.', '') # e.g., "pdf"

    input_filename = f"{original_name}{original_ext}"
    input_path = os.path.join(app.config['UPLOAD_FOLDER'], input_filename)
    file.save(input_path)

    output_filename = f"{original_name}_converted.{target_format}"
    output_path = os.path.join(app.config['UPLOAD_FOLDER'], output_filename)
    
    msg = ""
    
    # --- UPDATED Conversion Logic ---
    try:
        # PDF Conversions
        if source_ext == 'pdf':
            if target_format == 'docx':
                msg = pdf_to_docx(input_path, output_path)
            elif target_format == 'pptx':
                msg = pdf_to_ppt(input_path, output_path)
            elif target_format == 'txt':
                msg = pdf_to_txt(input_path, output_path)
        
        # Image Conversions
        elif source_ext in ['jpg', 'jpeg', 'png', 'webp']:
            if target_format == 'pdf':
                msg = image_to_pdf(input_path, output_path)
            elif target_format in ['jpg', 'jpeg', 'png', 'webp']:
                msg = image_to_image(input_path, output_path, target_format)

        else:
            msg = f"Unsupported conversion: From {source_ext} to {target_format}"
    
    except Exception as e:
        msg = f"A server error occurred: {str(e)}"
    # -------------------------------

    # Clean up the INPUT file
    if os.path.exists(input_path):
        os.remove(input_path)

    # Send file if successful
    if "Converted successfully" in msg and os.path.exists(output_path):
        g.cleanup_file = output_path
        return send_file(output_path, as_attachment=True, download_name=output_filename)
    
    else:
        # Clean up the OUTPUT file if conversion failed
        if os.path.exists(output_path):
            os.remove(output_path)
        return f"<h3>Conversion Failed: {msg}</h3><a href='/'>Go Back</a>", 500


@app.after_request
def cleanup_files(response):
    """Cleans up the sent file after the request is finished."""
    file_to_delete = getattr(g, 'cleanup_file', None)
    if file_to_delete and os.path.exists(file_to_delete):
        try:
            os.remove(file_to_delete)
        except Exception as e:
            app.logger.error(f"Error cleaning up file {file_to_delete}: {e}")
    return response


if __name__ == "__main__":
    app.run(host='0.0.0.0', port=5000, debug=True)